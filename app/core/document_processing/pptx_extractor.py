"""PowerPoint document extractor with vision analysis for image-heavy slides.

Uses python-pptx for native text extraction from slides, tables, and speaker notes.
Falls back to Claude Vision (Haiku) for slides that are primarily images.
"""

import base64
import io
from typing import Any

from app.core.document_processing.base import (
    BaseExtractor,
    DocumentType,
    ExtractedSection,
    ExtractionError,
    ExtractionResult,
    ExtractorRegistry,
    PAGE_LIMITS,
    SIZE_LIMITS,
)
from app.core.logging import get_logger

logger = get_logger(__name__)

# Lazy imports
pptx_module = None


def _get_pptx():
    """Lazy load python-pptx."""
    global pptx_module
    if pptx_module is None:
        try:
            import pptx as _pptx

            pptx_module = _pptx
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX extraction. "
                "Install with: pip install python-pptx"
            )
    return pptx_module


# Minimum text length to consider a slide as "text-based"
MIN_TEXT_CHARS = 100
# Skip tiny images (icons, bullets)
MIN_IMAGE_BYTES = 5 * 1024  # 5KB
# Max vision API calls per document
MAX_VISION_CALLS = 20
# Max images per vision API call
MAX_IMAGES_PER_CALL = 4

# Vision prompt for image-heavy slides
SLIDE_VISION_PROMPT = """Analyze this image from a PowerPoint slide and extract all relevant content for a software requirements project.

Your task:
1. Extract ALL visible text exactly as shown
2. Describe diagrams, charts, flowcharts, or visual elements
3. Note any UI components, wireframes, or mockups
4. Identify data, metrics, or technical details shown

Return your analysis as structured text. Be thorough but concise. Focus on extractable information useful for software requirements."""


class PPTXExtractor(BaseExtractor):
    """PowerPoint document extractor.

    Extracts text from PPTX presentations including:
    - Slide titles and body text
    - Tables as markdown
    - Speaker notes
    - Image-heavy slides via Claude Vision (Haiku)
    """

    def can_handle(self, mime_type: str, file_extension: str) -> bool:
        """Check if this extractor can handle the file."""
        return mime_type in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ) or file_extension.lower() in (".pptx", ".ppt")

    def get_supported_types(self) -> list[str]:
        """Return supported MIME types."""
        return [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ]

    def get_size_limit(self) -> int:
        """Get size limit for PPTX."""
        return SIZE_LIMITS.get(DocumentType.PPTX, 15 * 1024 * 1024)

    async def extract(
        self,
        file_bytes: bytes,
        filename: str,
        max_pages: int | None = None,
        **kwargs: Any,
    ) -> ExtractionResult:
        """Extract content from PPTX.

        Args:
            file_bytes: Raw PPTX content
            filename: Original filename
            max_pages: Override max slides (default from PAGE_LIMITS)
            **kwargs: Additional options (extract_images=True to collect embedded images)

        Returns:
            ExtractionResult with extracted sections

        Raises:
            ExtractionError: If extraction fails
        """
        valid, error_msg = self.validate_size(file_bytes)
        if not valid:
            raise ExtractionError(error_msg, extractor="pptx", recoverable=False)

        pptx_lib = _get_pptx()
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        max_slides = max_pages or PAGE_LIMITS.get(DocumentType.PPTX, 50)

        try:
            pptx_stream = io.BytesIO(file_bytes)
            prs = pptx_lib.Presentation(pptx_stream)

            slide_count = len(prs.slides)
            warnings: list[str] = []

            if slide_count > max_slides:
                warnings.append(
                    f"Presentation has {slide_count} slides, truncating to {max_slides}"
                )
                slide_count = max_slides

            sections: list[ExtractedSection] = []
            all_text_parts: list[str] = []
            total_words = 0
            embedded_images: list[bytes] = []
            vision_calls = 0
            text_slides = 0
            image_slides = 0

            # Collect image-heavy slides for batch vision analysis
            image_slide_queue: list[dict] = []

            for slide_idx in range(slide_count):
                slide = prs.slides[slide_idx]
                slide_num = slide_idx + 1

                # Extract slide title
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()

                section_path = f"Slide {slide_num}"
                if title:
                    section_path = f"Slide {slide_num} > {title}"

                # Extract all text from shapes
                slide_text_parts: list[str] = []
                slide_images: list[bytes] = []

                for shape in slide.shapes:
                    # Text frames
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                # Preserve bullet level
                                indent = "  " * paragraph.level if paragraph.level else ""
                                prefix = "- " if paragraph.level > 0 else ""
                                slide_text_parts.append(f"{indent}{prefix}{text}")

                    # Tables
                    if shape.has_table:
                        table_md = self._extract_table(shape.table)
                        if table_md:
                            slide_text_parts.append(table_md)

                    # Images
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            img_blob = shape.image.blob
                            if len(img_blob) >= MIN_IMAGE_BYTES:
                                slide_images.append(img_blob)
                                embedded_images.append(img_blob)
                        except Exception:
                            pass

                # Extract speaker notes
                notes_text = ""
                try:
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        raw_notes = slide.notes_slide.notes_text_frame.text.strip()
                        if raw_notes:
                            notes_text = raw_notes
                except Exception:
                    pass

                slide_text = "\n".join(slide_text_parts)
                is_image_heavy = len(slide_text) < MIN_TEXT_CHARS and len(slide_images) > 0

                if is_image_heavy:
                    image_slides += 1
                    # Queue for vision analysis
                    image_slide_queue.append({
                        "slide_num": slide_num,
                        "title": title,
                        "section_path": section_path,
                        "images": slide_images,
                        "text": slide_text,
                        "notes": notes_text,
                    })
                else:
                    text_slides += 1

                # Always add whatever text we extracted
                if slide_text:
                    sections.append(
                        ExtractedSection(
                            section_type="paragraph",
                            content=slide_text,
                            section_title=title or f"Slide {slide_num}",
                            page_number=slide_num,
                            section_path=section_path,
                        )
                    )
                    all_text_parts.append(slide_text)
                    total_words += len(slide_text.split())

                # Add speaker notes as separate section
                if notes_text:
                    sections.append(
                        ExtractedSection(
                            section_type="speaker_notes",
                            content=notes_text,
                            section_title=f"Notes: {title or f'Slide {slide_num}'}",
                            page_number=slide_num,
                            section_path=f"{section_path} > Notes",
                        )
                    )
                    all_text_parts.append(f"[Speaker Notes] {notes_text}")
                    total_words += len(notes_text.split())

            # Process image-heavy slides with vision
            if image_slide_queue:
                vision_sections, vision_words, vision_calls = await self._process_image_slides(
                    image_slide_queue
                )
                sections.extend(vision_sections)
                total_words += vision_words
                for vs in vision_sections:
                    all_text_parts.append(vs.content)

            # Determine extraction method
            if image_slides > text_slides and vision_calls > 0:
                extraction_method = "vision"
            elif image_slides > 0 and vision_calls > 0:
                extraction_method = "hybrid"
            else:
                extraction_method = "native"

            logger.info(
                f"Extracted PPTX {filename}: {slide_count} slides, "
                f"{len(sections)} sections, {total_words} words, "
                f"method={extraction_method}, vision_calls={vision_calls}"
            )

            return ExtractionResult(
                sections=sections,
                page_count=slide_count,
                word_count=total_words,
                extraction_method=extraction_method,
                raw_text="\n\n".join(all_text_parts),
                embedded_images=embedded_images,
                metadata={
                    "filename": filename,
                    "text_slides": text_slides,
                    "image_slides": image_slides,
                    "vision_calls": vision_calls,
                    "total_images": len(embedded_images),
                },
                warnings=warnings,
            )

        except ExtractionError:
            raise
        except Exception as e:
            logger.error(f"PPTX extraction failed for {filename}: {e}")
            raise ExtractionError(
                f"PPTX extraction failed: {e}",
                extractor="pptx",
                recoverable=True,
            )

    def _extract_table(self, table) -> str:
        """Convert a PPTX table to markdown format."""
        rows = []
        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if row_idx == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")

        return "\n".join(rows) if rows else ""

    async def _process_image_slides(
        self,
        slide_queue: list[dict],
    ) -> tuple[list[ExtractedSection], int, int]:
        """Process image-heavy slides with Claude Vision.

        Batches up to MAX_IMAGES_PER_CALL images per API call.

        Returns:
            Tuple of (sections, total_words, vision_calls)
        """
        from anthropic import Anthropic

        from app.core.config import get_settings
        from app.core.llm_usage import log_llm_usage

        settings = get_settings()
        if not settings.ANTHROPIC_API_KEY:
            logger.warning("No ANTHROPIC_API_KEY, skipping vision analysis for image slides")
            return [], 0, 0

        client = Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        model = "claude-haiku-4-5-20251001"
        sections: list[ExtractedSection] = []
        total_words = 0
        vision_calls = 0

        for slide_info in slide_queue:
            if vision_calls >= MAX_VISION_CALLS:
                logger.warning("Hit max vision calls limit, skipping remaining image slides")
                break

            images = slide_info["images"][:MAX_IMAGES_PER_CALL]
            slide_num = slide_info["slide_num"]
            title = slide_info["title"]
            section_path = slide_info["section_path"]

            # Build message content with images
            content: list[dict] = []
            for img_bytes in images:
                mime_type = self._detect_image_mime(img_bytes)
                img_b64 = base64.standard_b64encode(img_bytes).decode("utf-8")
                content.append({
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": mime_type,
                        "data": img_b64,
                    },
                })

            content.append({
                "type": "text",
                "text": f"Slide {slide_num}: \"{title or 'Untitled'}\"\n\n{SLIDE_VISION_PROMPT}",
            })

            try:
                response = client.messages.create(
                    model=model,
                    max_tokens=2048,
                    messages=[{"role": "user", "content": content}],
                )

                vision_calls += 1
                analysis_text = response.content[0].text if response.content else ""

                # Log LLM usage
                usage = response.usage
                log_llm_usage(
                    workflow="document_processing",
                    model=model,
                    provider="anthropic",
                    tokens_input=usage.input_tokens,
                    tokens_output=usage.output_tokens,
                    chain="pptx_vision",
                )

                if analysis_text:
                    word_count = len(analysis_text.split())
                    total_words += word_count
                    sections.append(
                        ExtractedSection(
                            section_type="image_description",
                            content=analysis_text,
                            section_title=title or f"Slide {slide_num} (Vision)",
                            page_number=slide_num,
                            section_path=section_path,
                            metadata={"vision_analyzed": True, "model": model},
                        )
                    )

            except Exception as e:
                logger.warning(f"Vision analysis failed for slide {slide_num}: {e}")

        return sections, total_words, vision_calls

    @staticmethod
    def _detect_image_mime(img_bytes: bytes) -> str:
        """Detect image MIME type from magic bytes."""
        if img_bytes[:8] == b"\x89PNG\r\n\x1a\n":
            return "image/png"
        elif img_bytes[:2] == b"\xff\xd8":
            return "image/jpeg"
        elif img_bytes[:4] == b"RIFF" and img_bytes[8:12] == b"WEBP":
            return "image/webp"
        elif img_bytes[:6] in (b"GIF87a", b"GIF89a"):
            return "image/gif"
        # Default to PNG
        return "image/png"


# Register the extractor
ExtractorRegistry.register(PPTXExtractor())
